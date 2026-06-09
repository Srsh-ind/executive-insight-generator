from pptx import Presentation
from pptx.util import Inches


def create_ppt(insights, risks, recommendations, chart_path):
    prs = Presentation()

    # Slide 1 - Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Executive Insights"
    slide.placeholders[1].text = "AI-generated business performance summary"

    # Slide 2 - Key Insights
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Insights"
    slide.placeholders[1].text = "\n".join(insights)

    # Slide 3 - Revenue Chart
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Revenue by Region"
    slide.shapes.add_picture(chart_path, Inches(1), Inches(1.5), width=Inches(8))

    # Slide 4 - Risks
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Risks"
    slide.placeholders[1].text = "\n".join(risks)

    # Slide 5 - Recommendations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Recommendations"
    slide.placeholders[1].text = "\n".join(recommendations)

    ppt_file = "outputs/output.pptx"
    prs.save(ppt_file)

    return ppt_file